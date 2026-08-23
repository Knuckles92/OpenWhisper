"""Consent-gated, bounded search of a user-selected knowledge folder.

Meeting agents call this through ``AgentToolHost.search_context_files``.
Policy lives here so no caller can bypass consent, root containment, caps,
or the rule that live evidence ids never leak from folder text — they share
the ``sg_`` prefix of transcript evidence ids and would be fuzzy-repaired
onto the current meeting.
"""
from __future__ import annotations

import logging
import os
import re
import threading
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

logger = logging.getLogger(__name__)

DEFAULT_HIT_LIMIT = 10
MAX_HIT_LIMIT = 20
MAX_SNIPPET_CHARS = 180
MAX_PASSAGE_CHARS = 800
MAX_TOTAL_CHARS = 4000
MAX_FILE_BYTES = 2 * 1024 * 1024
MAX_DOC_BYTES = 8 * 1024 * 1024
MAX_EXTRACTED_CHARS = 50_000
MAX_FILES_SCANNED = 200
MAX_PDF_PAGES = 40
MAX_XLSX_ROWS = 200
MAX_XLSX_COLS = 20
MAX_CACHE_ENTRIES = 64

_SG_ID_RE = re.compile(r"\bsg_[0-9a-fA-F]{6,}\b")
_TOKEN_RE = re.compile(r"[A-Za-z0-9_]{2,}")

_SKIP_DIR_NAMES = frozenset({
    ".obsidian",
    ".git",
    ".svn",
    ".hg",
    ".trash",
    ".smart-env",
    ".idea",
    ".vscode",
    ".venv",
    "venv",
    "__pycache__",
    "node_modules",
})
_DOC_EXTS = frozenset({".pdf", ".docx", ".pptx", ".xlsx"})
_TEXT_EXTS = frozenset({
    ".md", ".markdown", ".mdx", ".txt", ".rst", ".csv", ".tsv",
    ".json", ".jsonl", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".conf", ".xml", ".html", ".htm", ".css", ".svg",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".go", ".rs", ".java",
    ".kt", ".c", ".h", ".cpp", ".hpp", ".rb", ".php", ".sh",
    ".sql", ".r", ".tex", ".log", ".env.example",
})
_SKIP_EXTS = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".ico",
    ".tif", ".tiff", ".heic", ".avif",
    ".mp3", ".wav", ".m4a", ".flac", ".ogg", ".aac",
    ".mp4", ".mov", ".mkv", ".avi", ".webm",
    ".zip", ".tar", ".gz", ".bz2", ".7z", ".rar",
    ".exe", ".dll", ".so", ".dylib", ".bin", ".o", ".obj",
    ".pyc", ".class", ".wasm",
    ".doc", ".ppt", ".xls", ".odt", ".odp", ".ods",
})

_CACHE_LOCK = threading.Lock()
_TEXT_CACHE: Dict[Tuple[str, str, int, int], str] = {}
_CACHE_ORDER: List[Tuple[str, str, int, int]] = []

__all__ = [
    "DEFAULT_HIT_LIMIT",
    "MAX_HIT_LIMIT",
    "MAX_SNIPPET_CHARS",
    "MAX_PASSAGE_CHARS",
    "MAX_TOTAL_CHARS",
    "search_context_files",
    "clear_context_folder_cache",
]


def search_context_files(
    query: str = "",
    relative_path: Optional[str] = None,
    limit: int = DEFAULT_HIT_LIMIT,
) -> Dict[str, Any]:
    """Bounded, consent-gated keyword search across the knowledge folder.

    Args:
        query: Free-text search terms. Required unless ``relative_path`` is set.
        relative_path: Optional file inside the folder for a bounded passage.
        limit: Requested hit cap. Clamped to ``MAX_HIT_LIMIT``.

    Returns:
        ``{"ok", "disabled"?, "text", "hits"}``. ``text`` is what the model
        sees. Hits never include ``sg_`` segment ids or absolute paths.
    """
    if not _folder_enabled():
        return _disabled("Knowledge-folder search is disabled.")
    root = _configured_root()
    if root is None:
        return _disabled("Knowledge-folder search is not available.")

    hit_limit = _clamp_limit(limit)
    target = str(relative_path or "").strip() or None

    try:
        if target:
            return _slice_file(root, relative_path=target, query=query)
        return _search_folder(root, query=query, limit=hit_limit)
    except Exception:
        logger.exception("Knowledge-folder search failed")
        return {
            "ok": False,
            "text": "Knowledge-folder search failed.",
            "hits": [],
        }


def clear_context_folder_cache() -> None:
    """Drop the in-memory extraction cache (tests and settings changes)."""
    with _CACHE_LOCK:
        _TEXT_CACHE.clear()
        _CACHE_ORDER.clear()


def _folder_enabled() -> bool:
    try:
        from services.settings import resolve_meeting_context_folder_enabled

        return bool(resolve_meeting_context_folder_enabled())
    except Exception:
        return False


def _configured_root() -> Optional[Path]:
    try:
        from services.settings import resolve_meeting_context_folder_path

        raw = resolve_meeting_context_folder_path()
    except Exception:
        return None
    cleaned = str(raw or "").strip()
    if not cleaned:
        return None
    try:
        root = Path(cleaned).expanduser().resolve(strict=True)
    except OSError:
        return None
    if not root.is_dir() or root.is_symlink():
        return None
    return root


def _clamp_limit(limit: Any) -> int:
    try:
        value = int(limit)
    except (TypeError, ValueError):
        value = DEFAULT_HIT_LIMIT
    return max(1, min(value, MAX_HIT_LIMIT))


def _disabled(text: str) -> Dict[str, Any]:
    return {"ok": False, "disabled": True, "text": text, "hits": []}


def _search_folder(
    root: Path, *, query: str, limit: int,
) -> Dict[str, Any]:
    tokens = _query_tokens(query)
    if not tokens:
        return {
            "ok": False,
            "text": "Provide a search query or a relative file path.",
            "hits": [],
        }

    scored: List[Tuple[int, str, str]] = []
    scanned = 0
    for path, relpath in _iter_files(root):
        scanned += 1
        text = _cached_text(root, path, relpath)
        if not text:
            continue
        score = _score(relpath, text, tokens)
        if score <= 0:
            continue
        scored.append((score, relpath, text))

    scored.sort(key=lambda item: (-item[0], item[1].lower()))
    hits: List[Dict[str, Any]] = []
    lines: List[str] = []
    budget = MAX_TOTAL_CHARS
    for index, (_item_score, relpath, text) in enumerate(scored[:limit], 1):
        snippet = _snippet(text, tokens, MAX_SNIPPET_CHARS)
        if not snippet:
            continue
        ref = _opaque_ref(relpath, index)
        line = f"{ref}  {relpath} — {snippet}"
        if len(line) > budget:
            break
        lines.append(line)
        budget -= len(line) + 1
        hits.append({
            "ref": ref,
            "path": relpath,
            "title": Path(relpath).name,
            "snippet": snippet,
        })
    if not hits:
        extra = ""
        if scanned >= MAX_FILES_SCANNED:
            extra = " Scan limit reached."
        return {
            "ok": True,
            "text": f"No knowledge-folder matches.{extra}".rstrip(),
            "hits": [],
        }
    return {"ok": True, "text": "\n".join(lines), "hits": hits}


def _slice_file(
    root: Path, *, relative_path: str, query: str,
) -> Dict[str, Any]:
    path = _safe_join(root, relative_path)
    if path is None:
        return {
            "ok": False,
            "text": "That file is outside the knowledge folder.",
            "hits": [],
        }
    if not path.is_file():
        return {"ok": False, "text": "Unknown knowledge-folder file.", "hits": []}
    relpath = _relative_display(root, path)
    if relpath is None:
        return {
            "ok": False,
            "text": "That file is outside the knowledge folder.",
            "hits": [],
        }
    text = _cached_text(root, path, relpath)
    if not text:
        return {
            "ok": True,
            "text": f"No extractable text in {relpath}.",
            "hits": [],
        }
    tokens = _query_tokens(query)
    passage = (
        _snippet(text, tokens, MAX_PASSAGE_CHARS)
        if tokens else _clip(text, MAX_PASSAGE_CHARS)
    )
    if not passage:
        return {
            "ok": True,
            "text": f"No extractable text in {relpath}.",
            "hits": [],
        }
    ref = _opaque_ref(relpath, 1)
    line = f"{ref}  {relpath}\n{passage}"
    return {
        "ok": True,
        "text": _clip(line, MAX_TOTAL_CHARS),
        "hits": [{
            "ref": ref,
            "path": relpath,
            "title": Path(relpath).name,
            "snippet": passage,
        }],
    }


def _iter_files(root: Path) -> Iterable[Tuple[Path, str]]:
    scanned = 0
    for dirpath, dirnames, filenames in os.walk(root, followlinks=False):
        if os.path.islink(dirpath):
            dirnames[:] = []
            continue
        dirnames[:] = sorted(
            name for name in dirnames if not _skip_dir(name)
        )
        for name in sorted(filenames):
            if scanned >= MAX_FILES_SCANNED:
                return
            if name.startswith("."):
                continue
            full = os.path.join(dirpath, name)
            if os.path.islink(full) or not os.path.isfile(full):
                continue
            path = Path(full)
            relpath = _relative_display(root, path)
            if relpath is None:
                continue
            scanned += 1
            yield path, relpath


def _skip_dir(name: str) -> bool:
    lowered = name.lower()
    return lowered.startswith(".") or lowered in _SKIP_DIR_NAMES


def _safe_join(root: Path, relative_path: str) -> Optional[Path]:
    raw = (relative_path or "").replace("\\", "/").strip()
    if not raw or raw.startswith("/") or raw.startswith("~"):
        return None
    parts = Path(raw).parts
    if not parts or any(part in ("", "..") for part in parts):
        return None
    if any(part.startswith(".") for part in parts):
        return None
    candidate = root.joinpath(*parts)
    if candidate.is_symlink():
        return None
    try:
        resolved = candidate.resolve(strict=False)
        resolved.relative_to(root)
    except (OSError, ValueError):
        return None
    try:
        if os.path.islink(candidate) or os.path.realpath(candidate) != str(resolved):
            return None
    except OSError:
        return None
    return resolved


def _relative_display(root: Path, path: Path) -> Optional[str]:
    try:
        resolved = path.resolve(strict=True)
        rel = resolved.relative_to(root)
    except (OSError, ValueError):
        return None
    if any(part.startswith(".") or part.lower() in _SKIP_DIR_NAMES
           for part in rel.parts):
        return None
    return rel.as_posix()


def _cached_text(root: Path, path: Path, relpath: str) -> str:
    try:
        stat = path.stat()
    except OSError:
        return ""
    key = (str(root), relpath, int(stat.st_size), int(stat.st_mtime_ns))
    with _CACHE_LOCK:
        cached = _TEXT_CACHE.get(key)
        if cached is not None:
            return cached
    text = _extract_text(path)
    if not text:
        return ""
    with _CACHE_LOCK:
        _TEXT_CACHE[key] = text
        _CACHE_ORDER.append(key)
        while len(_CACHE_ORDER) > MAX_CACHE_ENTRIES:
            old = _CACHE_ORDER.pop(0)
            _TEXT_CACHE.pop(old, None)
    return text


def _extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in _SKIP_EXTS:
        return ""
    try:
        size = path.stat().st_size
    except OSError:
        return ""
    limit = MAX_DOC_BYTES if suffix in _DOC_EXTS else MAX_FILE_BYTES
    if size <= 0 or size > limit:
        return ""
    try:
        if suffix == ".pdf":
            extracted = _extract_pdf(path)
        elif suffix == ".docx":
            extracted = _extract_docx(path)
        elif suffix == ".pptx":
            extracted = _extract_pptx(path)
        elif suffix == ".xlsx":
            extracted = _extract_xlsx(path)
        elif suffix in _TEXT_EXTS:
            extracted = _extract_plain(path)
        else:
            extracted = _extract_plain(path)
    except Exception:
        logger.debug("Knowledge-folder extract failed for %s", path, exc_info=True)
        return ""
    return _clip(_sanitize(extracted or ""), MAX_EXTRACTED_CHARS)


def _extract_plain(path: Path) -> str:
    data = path.read_bytes()
    if b"\x00" in data[:8192]:
        return ""
    for encoding in ("utf-8-sig", "utf-8"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    sample = data[:4096]
    printable = sum(1 for byte in sample if 32 <= byte <= 126 or byte in (9, 10, 13))
    if sample and printable / len(sample) < 0.85:
        return ""
    return data.decode("latin-1", errors="ignore")


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if getattr(reader, "is_encrypted", False):
        try:
            if reader.decrypt("") == 0:
                return ""
        except Exception:
            return ""
    pages = []
    for index, page in enumerate(reader.pages):
        if index >= MAX_PDF_PAGES:
            break
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)
    return "\n".join(pages)


def _extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    parts: List[str] = []
    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [
                (cell.text or "").strip()
                for cell in row.cells
                if (cell.text or "").strip()
            ]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text or "").strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts: List[str] = []
        for sheet in workbook.worksheets:
            parts.append(str(sheet.title or "Sheet"))
            for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                if row_index >= MAX_XLSX_ROWS:
                    break
                cells = [
                    str(value).strip()
                    for value in list(row)[:MAX_XLSX_COLS]
                    if value is not None and str(value).strip()
                ]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    finally:
        workbook.close()


def _query_tokens(query: str) -> List[str]:
    return [token.lower() for token in _TOKEN_RE.findall(query or "")]


def _score(relpath: str, text: str, tokens: List[str]) -> int:
    name = relpath.lower()
    haystack = text.lower()
    score = 0
    matched = 0
    for token in tokens:
        in_name = token in name
        count = haystack.count(token)
        if in_name:
            score += 8
            matched += 1
        if count:
            score += min(count, 10)
            matched += 1
    if matched == 0:
        return 0
    if all(token in name for token in tokens):
        score += 12
    return score


def _snippet(text: str, tokens: List[str], max_chars: int) -> str:
    cleaned = _sanitize(text)
    lowered = cleaned.lower()
    start = 0
    for token in tokens:
        index = lowered.find(token)
        if index >= 0:
            start = max(0, index - max(24, max_chars // 5))
            break
    window = cleaned[start:start + max_chars + 40]
    return _clip(window, max_chars)


def _opaque_ref(relpath: str, index: int) -> str:
    return f"file:{relpath}:{index}"


def _clip(text: Any, max_chars: int) -> str:
    cleaned = _sanitize(str(text or "")).strip()
    if len(cleaned) <= max_chars:
        return cleaned
    return cleaned[: max_chars - 1].rstrip() + "…"


def _sanitize(text: str) -> str:
    return _SG_ID_RE.sub("[id]", text)
